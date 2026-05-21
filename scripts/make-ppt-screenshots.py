"""
소싱킷 결제경로 PPT - 실제 스크린샷 버전
토스페이먼츠 정기결제(빌링) 서비스 심사 제출용
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
from pathlib import Path
import os

# 스크린샷 경로
SHOT_DIR = Path(r"C:\Users\USER\Documents\Claude\cord\trade-sourcing-app\scripts\screenshots")
OUT_PATH = r"C:\Users\USER\Documents\Claude\cord\trade-sourcing-app\소싱킷_결제경로_v4.pptx"

# 색상
ORANGE    = RGBColor(0xFF, 0x78, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1A, 0x1A, 0x2E)
GRAY      = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG  = RGBColor(0xF8, 0xF9, 0xFA)
TOSS_BLUE = RGBColor(0x00, 0x78, 0xFF)
GREEN     = RGBColor(0x10, 0xB9, 0x81)

# 슬라이드 크기 (16:9)
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def blank_slide():
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

def add_rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = lw if lw else Pt(1)
    else:
        shape.line.fill.background()
    shape.line.width = lw
    return shape

def add_text(slide, text, l, t, w, h, size=Pt(14), bold=False, color=WHITE,
             align=PP_ALIGN.CENTER, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def prepare_img(img_path, suffix, top_crop=0, bottom_crop=0):
    """
    Playwright로 캡처한 모바일 스크린샷 처리
    - 이미 모바일 뷰포트(780x1688px)로 캡처됨 → 크롭 불필요
    - 상단/하단만 선택적으로 자르기
    """
    img = Image.open(img_path)
    iw, ih = img.size
    if top_crop > 0 or bottom_crop > 0:
        cropped = img.crop((0, top_crop, iw, ih - bottom_crop if bottom_crop > 0 else ih))
    else:
        cropped = img
    tmp = str(SHOT_DIR / f"_tmp_{suffix}.png")
    cropped.save(tmp)
    return tmp

def header(slide, title, subtitle=""):
    """공통 헤더 배경"""
    add_rect(slide, 0, 0, W, Inches(1.1), fill=DARK)
    # 소싱킷 로고 텍스트
    add_text(slide, "소싱킷", Inches(0.3), Inches(0.15), Inches(2), Inches(0.5),
             size=Pt(20), bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
    # 타이틀
    add_text(slide, title, Inches(2.5), Inches(0.1), Inches(7), Inches(0.6),
             size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, Inches(2.5), Inches(0.65), Inches(7), Inches(0.35),
                 size=Pt(12), color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

def step_badge(slide, num, label, x, y):
    """단계 번호 배지"""
    add_rect(slide, x, y, Inches(0.4), Inches(0.4), fill=ORANGE)
    add_text(slide, str(num), x, y, Inches(0.4), Inches(0.4),
             size=Pt(14), bold=True, color=WHITE)
    add_text(slide, label, x + Inches(0.45), y + Inches(0.02), Inches(2), Inches(0.36),
             size=Pt(11), bold=True, color=DARK, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ═══════════════════════════════════════════════════════════════
slide1 = blank_slide()
add_rect(slide1, 0, 0, W, H, fill=DARK)

# 오렌지 그라디언트 바
add_rect(slide1, 0, 0, W, Inches(0.08), fill=ORANGE)

# 중앙 콘텐츠 박스
add_rect(slide1, Inches(3), Inches(1.5), Inches(7.33), Inches(4.5),
         fill=RGBColor(0x25, 0x25, 0x3D))

# 서비스명
add_text(slide1, "소싱킷", Inches(3.5), Inches(1.8), Inches(6.33), Inches(0.9),
         size=Pt(44), bold=True, color=ORANGE)
add_text(slide1, "무역 소싱 관리 플랫폼",
         Inches(3.5), Inches(2.65), Inches(6.33), Inches(0.5),
         size=Pt(18), color=WHITE)

# 구분선
add_rect(slide1, Inches(4.5), Inches(3.25), Inches(4.33), Inches(0.04), fill=ORANGE)

# 부제목
add_text(slide1, "결제 경로 안내서",
         Inches(3.5), Inches(3.4), Inches(6.33), Inches(0.55),
         size=Pt(22), bold=True, color=WHITE)
add_text(slide1, "토스페이먼츠 서비스 심사 제출용",
         Inches(3.5), Inches(3.95), Inches(6.33), Inches(0.4),
         size=Pt(13), color=RGBColor(0xAA, 0xAA, 0xAA))

# URL
add_text(slide1, "www.sourcing-kit.kr",
         Inches(3.5), Inches(4.5), Inches(6.33), Inches(0.4),
         size=Pt(13), color=TOSS_BLUE)

# 하단
add_rect(slide1, 0, H - Inches(0.6), W, Inches(0.6), fill=RGBColor(0x11, 0x11, 0x22))
add_text(slide1, "sourcing-kit.kr  |  무역 소싱 자동화 플랫폼  |  2026",
         0, H - Inches(0.55), W, Inches(0.5),
         size=Pt(10), color=GRAY)

print("슬라이드 1 (표지) 완료")


# ═══════════════════════════════════════════════════════════════
# 슬라이드 2: 결제 흐름 개요
# ═══════════════════════════════════════════════════════════════
slide2 = blank_slide()
add_rect(slide2, 0, 0, W, H, fill=LIGHT_BG)
header(slide2, "결제 경로 개요", "소싱킷 맛보기(단건) · Pro(정기결제) 결제 흐름")

# 5단계 흐름도
steps = [
    ("1", "서비스 접속", "www.sourcing-kit.kr"),
    ("2", "로그인", "카카오 / 구글 / 이메일"),
    ("3", "요금제 선택", "맛보기 / Pro 월·연간"),
    ("4", "결제창 실행", "토스페이먼츠 SDK"),
    ("5", "결제 완료", "구독 활성화"),
]

bw = Inches(2.0)
bh = Inches(1.8)
gap = Inches(0.3)
total_w = len(steps) * bw + (len(steps) - 1) * gap
start_x = (W - total_w) / 2
y = Inches(2.2)

step_colors = [DARK, DARK, DARK, TOSS_BLUE, GREEN]

for i, (num, title, sub) in enumerate(steps):
    bx = start_x + i * (bw + gap)
    add_rect(slide2, bx, y, bw, bh, fill=step_colors[i])

    # 번호 원
    add_rect(slide2, bx + Inches(0.75), y + Inches(0.15),
             Inches(0.5), Inches(0.5), fill=ORANGE)
    add_text(slide2, num, bx + Inches(0.75), y + Inches(0.15),
             Inches(0.5), Inches(0.5), size=Pt(16), bold=True, color=WHITE)

    add_text(slide2, title, bx, y + Inches(0.75), bw, Inches(0.55),
             size=Pt(13), bold=True, color=WHITE)
    add_text(slide2, sub, bx, y + Inches(1.3), bw, Inches(0.4),
             size=Pt(10), color=RGBColor(0xCC, 0xCC, 0xCC))

    # 화살표
    if i < len(steps) - 1:
        ax = bx + bw + Inches(0.05)
        add_text(slide2, "→", ax, y + Inches(0.75), gap, Inches(0.5),
                 size=Pt(18), bold=True, color=ORANGE)

# 설명 박스
desc_y = Inches(4.4)
add_rect(slide2, Inches(1), desc_y, W - Inches(2), Inches(2.4),
         fill=WHITE, line=RGBColor(0xE0, 0xE0, 0xE0), lw=Pt(1))

desc_lines = [
    ("맛보기 (단건 결제)", "₩9,900 / 30일  |  requestPayment() → 카드 결제 → 결제 완료"),
    ("Pro 월간 (정기결제)", "₩7,900 / 월  |  requestBillingAuth() → 카드 등록 → 자동 청구"),
    ("Pro 연간 (정기결제)", "₩71,100 / 연 (25% 할인)  |  requestBillingAuth() → 카드 등록 → 자동 청구"),
]
for j, (lbl, desc) in enumerate(desc_lines):
    ry = desc_y + Inches(0.25) + j * Inches(0.6)
    add_rect(slide2, Inches(1.2), ry, Inches(2.2), Inches(0.35), fill=ORANGE)
    add_text(slide2, lbl, Inches(1.2), ry, Inches(2.2), Inches(0.35),
             size=Pt(10), bold=True, color=WHITE)
    add_text(slide2, desc, Inches(3.6), ry, Inches(8.5), Inches(0.35),
             size=Pt(10), color=DARK, align=PP_ALIGN.LEFT)

print("슬라이드 2 (개요) 완료")


# ═══════════════════════════════════════════════════════════════
# 슬라이드 3: 실제 화면 - 로그인
# ═══════════════════════════════════════════════════════════════
slide3 = blank_slide()
add_rect(slide3, 0, 0, W, H, fill=LIGHT_BG)
header(slide3, "STEP 1-2  |  서비스 접속 및 로그인")

# 왼쪽: 설명
lx = Inches(0.5)
ly = Inches(1.3)
lw = Inches(5.5)

step_badge(slide3, 1, "서비스 접속", lx, ly)
add_text(slide3,
         "www.sourcing-kit.kr 접속\n무역 소싱 관리 플랫폼",
         lx + Inches(0.45), ly + Inches(0.5), lw, Inches(0.8),
         size=Pt(11), color=DARK, align=PP_ALIGN.LEFT)

step_badge(slide3, 2, "로그인 선택", lx, ly + Inches(1.5))
login_desc = (
    "카카오로 시작하기\n"
    "구글로 시작하기\n"
    "이메일로 로그인"
)
add_text(slide3, login_desc,
         lx + Inches(0.45), ly + Inches(2.0), lw, Inches(1.2),
         size=Pt(11), color=DARK, align=PP_ALIGN.LEFT)

# 안내 박스
add_rect(slide3, lx, ly + Inches(3.5), lw, Inches(0.8),
         fill=RGBColor(0xFF, 0xF3, 0xE0), line=ORANGE, lw=Pt(1))
add_text(slide3, "소셜 로그인(카카오, 구글) 및 이메일 로그인 지원\n"
         "NextAuth 기반 세션 관리",
         lx, ly + Inches(3.5), lw, Inches(0.8),
         size=Pt(10), color=DARK, align=PP_ALIGN.CENTER)

# 오른쪽: 실제 스크린샷
img_path = str(SHOT_DIR / "01_login.png")
if Path(img_path).exists():
    tmp = prepare_img(img_path, "login", top_crop=0, bottom_crop=0)
    pw = Inches(2.8)  # 너비 고정, 높이 자동 (비율 유지)
    px = Inches(6.5)
    py = Inches(1.1)
    slide3.shapes.add_picture(tmp, px, py, pw)  # height auto
    pic = slide3.shapes[-1]
    ph = pic.height
    add_rect(slide3, px - Inches(0.05), py - Inches(0.05),
             pw + Inches(0.1), ph + Inches(0.1), line=DARK, lw=Pt(2))
    add_text(slide3, "실제 화면", px, py - Inches(0.35), pw, Inches(0.3),
             size=Pt(10), color=GRAY, bold=True)

print("슬라이드 3 (로그인) 완료")


# ═══════════════════════════════════════════════════════════════
# 슬라이드 4: 실제 화면 - 요금제 선택
# ═══════════════════════════════════════════════════════════════
slide4 = blank_slide()
add_rect(slide4, 0, 0, W, H, fill=LIGHT_BG)
header(slide4, "STEP 3  |  요금제 선택")

# 왼쪽 설명
lx = Inches(0.5)
ly = Inches(1.3)
lw = Inches(5.5)

step_badge(slide4, 3, "요금제 선택", lx, ly)

plans = [
    ("무료", "₩0", "상품 10개, 공급업체 5개"),
    ("맛보기", "₩9,900/30일", "상품 100개, 공급업체 30개"),
    ("Pro 월간", "₩7,900/월", "무제한 + 팀 공유"),
    ("Pro 연간", "₩71,100/연", "25% 할인"),
]
for j, (name, price, feat) in enumerate(plans):
    ry = ly + Inches(0.7) + j * Inches(0.85)
    bc = ORANGE if "맛보기" in name or "Pro" in name else DARK
    add_rect(slide4, lx, ry, Inches(1.2), Inches(0.6), fill=bc)
    add_text(slide4, name, lx, ry, Inches(1.2), Inches(0.6),
             size=Pt(10), bold=True, color=WHITE)
    add_text(slide4, price, lx + Inches(1.3), ry + Inches(0.02),
             Inches(1.8), Inches(0.3), size=Pt(11), bold=True, color=DARK,
             align=PP_ALIGN.LEFT)
    add_text(slide4, feat, lx + Inches(1.3), ry + Inches(0.3),
             Inches(4.0), Inches(0.3), size=Pt(10), color=GRAY, align=PP_ALIGN.LEFT)

add_rect(slide4, lx, ly + Inches(4.25), lw, Inches(0.75),
         fill=RGBColor(0xE8, 0xF5, 0xE9), line=GREEN, lw=Pt(1))
add_text(slide4, "요금제 선택 후 결제 버튼 클릭 → 토스페이먼츠 결제창 실행",
         lx, ly + Inches(4.25), lw, Inches(0.75),
         size=Pt(10), color=DARK, align=PP_ALIGN.CENTER)

# 오른쪽: 실제 스크린샷
img_path = str(SHOT_DIR / "03_pricing.png")
if Path(img_path).exists():
    tmp = prepare_img(img_path, "pricing", top_crop=0, bottom_crop=0)
    pw = Inches(2.8)
    px = Inches(6.5)
    py = Inches(1.1)
    slide4.shapes.add_picture(tmp, px, py, pw)
    pic = slide4.shapes[-1]
    ph = pic.height
    add_rect(slide4, px - Inches(0.05), py - Inches(0.05),
             pw + Inches(0.1), ph + Inches(0.1), line=DARK, lw=Pt(2))
    add_text(slide4, "실제 화면", px, py - Inches(0.35), pw, Inches(0.3),
             size=Pt(10), color=GRAY, bold=True)

print("슬라이드 4 (요금제) 완료")


# ═══════════════════════════════════════════════════════════════
# 슬라이드 5: 실제 화면 - 토스 결제창
# ═══════════════════════════════════════════════════════════════
slide5 = blank_slide()
add_rect(slide5, 0, 0, W, H, fill=LIGHT_BG)
header(slide5, "STEP 4  |  토스페이먼츠 결제창",
       "단건 결제: requestPayment() / 정기결제: requestBillingAuth()")

lx = Inches(0.5)
ly = Inches(1.3)
lw = Inches(5.5)

step_badge(slide5, 4, "결제창 실행", lx, ly)

# API 코드 박스
add_rect(slide5, lx, ly + Inches(0.6), lw, Inches(2.0),
         fill=DARK, line=RGBColor(0x44, 0x44, 0x66), lw=Pt(1))
code_text = (
    "// 맛보기 (단건 결제)\n"
    "await payment.requestPayment({\n"
    "  method: 'CARD',\n"
    "  amount: { currency: 'KRW', value: 9900 },\n"
    "  orderName: '소싱킷 맛보기 30일'\n"
    "})\n\n"
    "// Pro (정기결제 카드 등록)\n"
    "await payment.requestBillingAuth({\n"
    "  method: 'CARD'\n"
    "})"
)
add_text(slide5, code_text, lx + Inches(0.2), ly + Inches(0.7),
         lw - Inches(0.4), Inches(1.85),
         size=Pt(9), color=RGBColor(0x7D, 0xD3, 0xFC), align=PP_ALIGN.LEFT)

# 결제 방법 목록
add_rect(slide5, lx, ly + Inches(2.85), lw, Inches(1.2),
         fill=WHITE, line=RGBColor(0xE0, 0xE0, 0xE0), lw=Pt(1))
add_text(slide5, "결제 수단",
         lx + Inches(0.15), ly + Inches(2.9), lw, Inches(0.35),
         size=Pt(11), bold=True, color=DARK, align=PP_ALIGN.LEFT)
methods = "카카오페이  |  신한카드  |  하나Pay  |  우리카드\n케이뱅크  |  토스뱅크  |  카카오뱅크"
add_text(slide5, methods,
         lx + Inches(0.15), ly + Inches(3.2), lw, Inches(0.65),
         size=Pt(10), color=GRAY, align=PP_ALIGN.LEFT)

add_rect(slide5, lx, ly + Inches(4.2), lw, Inches(0.6),
         fill=TOSS_BLUE)
add_text(slide5, "Powered by toss payments",
         lx, ly + Inches(4.2), lw, Inches(0.6),
         size=Pt(12), bold=True, color=WHITE)

# 오른쪽: 단건결제 + 빌링결제 두 화면 나란히
py = Inches(1.1)
pw_each = Inches(2.2)

# 단건결제 (맛보기)
img_path1 = str(SHOT_DIR / "04_toss_onetime.png")
if Path(img_path1).exists():
    tmp1 = prepare_img(img_path1, "onetime")
    px1 = Inches(6.2)
    slide5.shapes.add_picture(tmp1, px1, py, pw_each)
    pic1 = slide5.shapes[-1]
    ph1 = pic1.height
    add_rect(slide5, px1 - Inches(0.05), py - Inches(0.05),
             pw_each + Inches(0.1), ph1 + Inches(0.1), line=TOSS_BLUE, lw=Pt(2))
    add_text(slide5, "단건 결제 화면", px1, py - Inches(0.35), pw_each, Inches(0.3),
             size=Pt(9), color=GRAY, bold=True)

# 빌링결제 (정기)
img_path2 = str(SHOT_DIR / "05_toss_billing.png")
if Path(img_path2).exists():
    tmp2 = prepare_img(img_path2, "billing")
    px2 = Inches(8.8)
    slide5.shapes.add_picture(tmp2, px2, py, pw_each)
    pic2 = slide5.shapes[-1]
    ph2 = pic2.height
    add_rect(slide5, px2 - Inches(0.05), py - Inches(0.05),
             pw_each + Inches(0.1), ph2 + Inches(0.1), line=ORANGE, lw=Pt(2))
    add_text(slide5, "정기결제 카드 등록", px2, py - Inches(0.35), pw_each, Inches(0.3),
             size=Pt(9), color=GRAY, bold=True)

print("슬라이드 5 (결제창) 완료")


# ═══════════════════════════════════════════════════════════════
# 슬라이드 6: 결제 완료
# ═══════════════════════════════════════════════════════════════
slide6 = blank_slide()
add_rect(slide6, 0, 0, W, H, fill=LIGHT_BG)
header(slide6, "STEP 5  |  결제 완료 및 구독 활성화")

lx = Inches(0.5)
ly = Inches(1.3)
lw = Inches(5.5)

step_badge(slide6, 5, "결제 완료", lx, ly)

# 완료 후 처리 설명
flow_items = [
    (GREEN,    "결제 성공",     "/payment/billing?planType=taste 리다이렉트"),
    (DARK,     "빌링키 발급",   "issueBillingKey() → DB 저장"),
    (DARK,     "구독 활성화",   "tenant.planType 업데이트"),
    (ORANGE,   "홈 화면 이동",  "플랜 혜택 즉시 적용"),
]

for j, (color, title, desc) in enumerate(flow_items):
    ry = ly + Inches(0.6) + j * Inches(0.9)
    add_rect(slide6, lx, ry, Inches(1.5), Inches(0.55), fill=color)
    add_text(slide6, title, lx, ry, Inches(1.5), Inches(0.55),
             size=Pt(11), bold=True, color=WHITE)
    add_text(slide6, desc,
             lx + Inches(1.7), ry + Inches(0.1), Inches(3.8), Inches(0.35),
             size=Pt(10), color=DARK, align=PP_ALIGN.LEFT)

# 서버 처리 박스
add_rect(slide6, lx, ly + Inches(4.1), lw, Inches(1.1),
         fill=DARK, line=GREEN, lw=Pt(1))
add_text(slide6,
         "POST /api/payments/billing\n"
         "confirmPayment() → issueBillingKey() → DB 업데이트",
         lx + Inches(0.2), ly + Inches(4.15), lw - Inches(0.4), Inches(1.0),
         size=Pt(10), color=RGBColor(0x7D, 0xD3, 0xFC), align=PP_ALIGN.LEFT)

# 오른쪽: 결제 완료 화면 (홈 화면으로 대체)
img_path = str(SHOT_DIR / "02_home.png")
if Path(img_path).exists():
    tmp = prepare_img(img_path, "home")
    pw = Inches(2.8)
    px = Inches(6.5)
    py = Inches(1.1)
    slide6.shapes.add_picture(tmp, px, py, pw)
    pic = slide6.shapes[-1]
    ph = pic.height
    add_rect(slide6, px - Inches(0.05), py - Inches(0.05),
             pw + Inches(0.1), ph + Inches(0.1), line=GREEN, lw=Pt(2))
    add_text(slide6, "결제 후 홈 화면", px, py - Inches(0.35), pw, Inches(0.3),
             size=Pt(10), color=GRAY, bold=True)

print("슬라이드 6 (완료) 완료")


# ═══════════════════════════════════════════════════════════════
# 슬라이드 7: 서비스 소개 (부가)
# ═══════════════════════════════════════════════════════════════
slide7 = blank_slide()
add_rect(slide7, 0, 0, W, H, fill=LIGHT_BG)
header(slide7, "소싱킷 서비스 소개")

features = [
    ("상품 소싱", "글로벌 공급업체 발굴\n원가계산 자동화"),
    ("실시간 환율", "한국은행 API 연동\n자동 환율 적용"),
    ("바이어 제안서", "AI 기반 제안서 생성\nPDF 내보내기"),
    ("시장 조사", "트렌드 분석\n경쟁 상품 조사"),
]

bw = Inches(2.8)
bh = Inches(2.2)
gap = Inches(0.4)

for i, (title, desc) in enumerate(features):
    row = i // 2
    col = i % 2
    bx = Inches(1.5) + col * (bw + gap)
    by = Inches(1.4) + row * (bh + gap)

    add_rect(slide7, bx, by, bw, bh, fill=WHITE,
             line=RGBColor(0xE0, 0xE0, 0xE0), lw=Pt(1))
    add_rect(slide7, bx, by, bw, Inches(0.45), fill=DARK)
    add_text(slide7, title, bx, by, bw, Inches(0.45),
             size=Pt(13), bold=True, color=WHITE)
    add_text(slide7, desc, bx, by + Inches(0.55), bw, Inches(1.4),
             size=Pt(11), color=DARK)

# 우측 - 수치
rx = Inches(8.5)
add_rect(slide7, rx, Inches(1.4), Inches(4), Inches(5.2),
         fill=DARK, line=ORANGE, lw=Pt(2))
add_text(slide7, "서비스 현황",
         rx, Inches(1.6), Inches(4), Inches(0.5),
         size=Pt(16), bold=True, color=ORANGE)

stats = [
    ("www.sourcing-kit.kr", "서비스 URL"),
    ("무역 전문 플랫폼", "서비스 카테고리"),
    ("맛보기 / Pro", "구독 플랜"),
    ("카카오·구글 로그인", "인증 방식"),
    ("한국 (KRW)", "결제 통화"),
]
for j, (val, lbl) in enumerate(stats):
    sy = Inches(2.2) + j * Inches(0.75)
    add_text(slide7, lbl, rx + Inches(0.2), sy, Inches(3.6), Inches(0.3),
             size=Pt(10), color=GRAY, align=PP_ALIGN.LEFT)
    add_text(slide7, val, rx + Inches(0.2), sy + Inches(0.28),
             Inches(3.6), Inches(0.35),
             size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

print("슬라이드 7 (소개) 완료")


# ═══════════════════════════════════════════════════════════════
# 슬라이드 8: 서비스 제공기간 및 환불정책 (Toss 심사 필수)
# ═══════════════════════════════════════════════════════════════
slide8 = blank_slide()
add_rect(slide8, 0, 0, W, H, fill=LIGHT_BG)
header(slide8, "서비스 제공기간 및 환불정책",
       "소싱킷 요금제 페이지 명시 내용 (www.sourcing-kit.kr/pricing)")

# ── 왼쪽: 서비스 제공기간 ──────────────────────────────────────
lx = Inches(0.5)
ly = Inches(1.2)
lw = Inches(5.8)

# 제목 배지
add_rect(slide8, lx, ly, lw, Inches(0.45), fill=DARK)
add_text(slide8, "📅  서비스 제공기간",
         lx, ly, lw, Inches(0.45),
         size=Pt(14), bold=True, color=WHITE)

period_rows = [
    ("무료 플랜",    "회원 탈퇴 시까지 무기한 제공"),
    ("맛보기",       "결제 완료일로부터 30일"),
    ("Pro 월구독",   "결제일로부터 다음 결제일 전날까지 (매월 자동 갱신)"),
    ("Pro 연구독",   "결제 완료일로부터 365일 (1년)"),
]

for j, (plan, desc) in enumerate(period_rows):
    ry = ly + Inches(0.55) + j * Inches(0.8)
    bc = ORANGE if "맛보기" in plan or "Pro" in plan else DARK
    add_rect(slide8, lx, ry, Inches(1.4), Inches(0.58), fill=bc)
    add_text(slide8, plan, lx, ry, Inches(1.4), Inches(0.58),
             size=Pt(10), bold=True, color=WHITE)
    add_rect(slide8, lx + Inches(1.5), ry, Inches(4.2), Inches(0.58),
             fill=WHITE, line=RGBColor(0xE0, 0xE0, 0xE0), lw=Pt(1))
    add_text(slide8, desc,
             lx + Inches(1.65), ry + Inches(0.08), Inches(4.0), Inches(0.45),
             size=Pt(10), color=DARK, align=PP_ALIGN.LEFT)

# ── 오른쪽: 환불 및 취소 정책 ─────────────────────────────────
rx = Inches(6.9)
rw = Inches(6.0)

add_rect(slide8, rx, ly, rw, Inches(0.45), fill=TOSS_BLUE)
add_text(slide8, "💳  환불 및 취소 정책",
         rx, ly, rw, Inches(0.45),
         size=Pt(14), bold=True, color=WHITE)

# 맛보기 섹션
ry_taste = ly + Inches(0.55)
add_rect(slide8, rx, ry_taste, rw, Inches(0.35), fill=RGBColor(0xFF, 0xF3, 0xE0))
add_text(slide8, "맛보기 (1회 단건결제)",
         rx + Inches(0.15), ry_taste, rw, Inches(0.35),
         size=Pt(11), bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

taste_items = [
    "• 서비스 이용 시작 전 취소 시 → 전액 환불",
    "• 서비스 이용 시작 후 → 환불 불가 (디지털 콘텐츠, 전자상거래법 제17조 제2항 제5호)",
    "• 서비스 결함·오류로 인한 경우 → 전액 환불",
]
for k, item in enumerate(taste_items):
    iy = ry_taste + Inches(0.4) + k * Inches(0.42)
    add_text(slide8, item,
             rx + Inches(0.15), iy, rw - Inches(0.2), Inches(0.38),
             size=Pt(10), color=DARK, align=PP_ALIGN.LEFT)

# Pro 구독 섹션
ry_pro = ry_taste + Inches(2.0)
add_rect(slide8, rx, ry_pro, rw, Inches(0.35), fill=RGBColor(0xE8, 0xF5, 0xE9))
add_text(slide8, "Pro 구독 (정기결제)",
         rx + Inches(0.15), ry_pro, rw, Inches(0.35),
         size=Pt(11), bold=True, color=GREEN, align=PP_ALIGN.LEFT)

pro_items = [
    "• 구독 취소는 언제든지 가능 (즉시 해지)",
    "• 취소 후 해당 결제 기간 종료일까지 서비스 이용 가능",
    "• 이미 결제된 구독료는 잔여 기간에 관계없이 환불 불가",
    "• 서비스 결함·오류로 인한 경우 → 잔여 기간 일할 환불",
]
for k, item in enumerate(pro_items):
    iy = ry_pro + Inches(0.4) + k * Inches(0.42)
    add_text(slide8, item,
             rx + Inches(0.15), iy, rw - Inches(0.2), Inches(0.38),
             size=Pt(10), color=DARK, align=PP_ALIGN.LEFT)

# 환불 연락처
contact_y = ly + Inches(5.5)
add_rect(slide8, lx, contact_y, W - Inches(1.0), Inches(0.6),
         fill=DARK, line=ORANGE, lw=Pt(1))
add_text(slide8,
         "환불 요청: goohw593@gmail.com  |  010-2623-6907  |  처리 기간: 영업일 기준 3~5일",
         lx, contact_y, W - Inches(1.0), Inches(0.6),
         size=Pt(11), bold=True, color=WHITE)

print("슬라이드 8 (환불정책) 완료")


# ═══════════════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════════════
prs.save(OUT_PATH)
print(f"\nPPT 저장 완료: {OUT_PATH}")
print(f"총 {len(prs.slides)}개 슬라이드")
