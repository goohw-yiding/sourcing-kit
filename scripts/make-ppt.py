"""
소싱킷 결제경로 PPT 생성 스크립트
토스페이먼츠 서비스 심사 제출용
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── 색상 팔레트 ──────────────────────────────────────────────
ORANGE      = RGBColor(0xFF, 0x78, 0x00)   # 소싱킷 오렌지
ORANGE_DARK = RGBColor(0xCC, 0x44, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x1A, 0x1A, 0x2E)
GRAY        = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY  = RGBColor(0xF3, 0xF4, 0xF6)
TOSS_BLUE   = RGBColor(0x00, 0x78, 0xFF)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
ARROW_RED   = RGBColor(0xEF, 0x44, 0x44)

# ── 슬라이드 크기 (16:9) ────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── 헬퍼 함수들 ─────────────────────────────────────────────
def blank_slide(prs):
    layout = prs.slide_layouts[6]  # 빈 레이아웃
    return prs.slides.add_slide(layout)

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width else Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=Pt(16), bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_img_placeholder(slide, l, t, w, h, label, bg_color=LIGHT_GRAY):
    """이미지 자리 표시자 (실제 스크린샷 대신 라벨로 표현)"""
    rect = add_rect(slide, l, t, w, h, fill_color=bg_color, line_color=GRAY, line_width=Pt(1.5))
    add_text(slide, label, l, t + h//2 - Pt(20), w, Pt(40),
             font_size=Pt(13), color=GRAY, align=PP_ALIGN.CENTER)
    return rect

def step_badge(slide, l, t, num, label):
    """번호 배지 + 텍스트"""
    # 원형 배지
    circle = slide.shapes.add_shape(9, l, t, Inches(0.45), Inches(0.45))  # 9=OVAL
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = str(num)
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    # 라벨
    add_text(slide, label, l + Inches(0.55), t, Inches(2.5), Inches(0.45),
             font_size=Pt(12), bold=True, color=DARK)

def arrow_right(slide, l, t):
    """→ 화살표"""
    add_text(slide, "→", l, t, Inches(0.4), Inches(0.4),
             font_size=Pt(20), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

def add_phone_frame(slide, l, t, w, h):
    """스마트폰 프레임"""
    # 외곽
    outer = add_rect(slide, l, t, w, h, fill_color=DARK, line_color=DARK, line_width=Pt(0))
    outer_shape = slide.shapes[-1]
    outer_shape.adjustments[0] = 0.06  # 라운드 코너 (지원 시)
    # 화면 영역
    margin = Inches(0.1)
    add_rect(slide, l + margin, t + Inches(0.25), w - margin*2, h - Inches(0.5),
             fill_color=WHITE, line_color=None)
    return outer

# ════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ════════════════════════════════════════════════════════════
slide1 = blank_slide(prs)

# 배경 그라디언트 효과 (좌측 오렌지)
add_rect(slide1, 0, 0, W, H, fill_color=ORANGE_DARK)
add_rect(slide1, Inches(7), 0, W - Inches(7), H, fill_color=RGBColor(0xFF, 0x99, 0x33))

# 제목
add_text(slide1, "소싱킷", Inches(1), Inches(1.5), Inches(6), Inches(1.2),
         font_size=Pt(52), bold=True, color=WHITE)
add_text(slide1, "결제 경로 안내서", Inches(1), Inches(2.8), Inches(6), Inches(0.8),
         font_size=Pt(28), bold=False, color=RGBColor(0xFF, 0xE0, 0xB0))
add_text(slide1, "무역 소싱 관리 서비스", Inches(1), Inches(3.7), Inches(6), Inches(0.6),
         font_size=Pt(18), color=RGBColor(0xFF, 0xCC, 0x99))

# 구분선
add_rect(slide1, Inches(1), Inches(4.4), Inches(4), Inches(0.04), fill_color=RGBColor(0xFF, 0xCC, 0x99))

# 사업자 정보
info_lines = [
    "이딩컴퍼니  |  대표자: 구희완",
    "사업자등록번호: 210-29-50637",
    "서비스 URL: https://sourcing-kit.kr",
    "결제 대행사: 토스페이먼츠(주)",
]
for i, line in enumerate(info_lines):
    add_text(slide1, line, Inches(1), Inches(4.6) + Pt(22)*i, Inches(7), Pt(22),
             font_size=Pt(12), color=RGBColor(0xFF, 0xE8, 0xCC))

# 우측 장식
add_text(slide1, "🔍", Inches(9.5), Inches(2), Inches(2), Inches(2),
         font_size=Pt(80), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 슬라이드 2: 전체 결제 흐름 요약
# ════════════════════════════════════════════════════════════
slide2 = blank_slide(prs)

# 헤더
add_rect(slide2, 0, 0, W, Inches(1.1), fill_color=ORANGE)
add_text(slide2, "결제 경로 전체 흐름", Inches(0.5), Inches(0.2), W - Inches(1), Inches(0.7),
         font_size=Pt(24), bold=True, color=WHITE)

# 5단계 박스
steps = [
    ("1", "서비스 접속",    "sourcing-kit.kr\n접속"),
    ("2", "플랜 선택",      "가격 페이지에서\n플랜 선택"),
    ("3", "결제하기 클릭",  "결제하기 버튼\n클릭"),
    ("4", "토스 결제창",    "토스페이먼츠\n결제 진행"),
    ("5", "결제 완료",      "서비스\n이용 시작"),
]

box_w = Inches(2.1)
box_h = Inches(3.5)
start_x = Inches(0.4)
top_y = Inches(1.5)
gap = Inches(0.25)

for i, (num, title, desc) in enumerate(steps):
    x = start_x + (box_w + gap) * i

    # 박스
    add_rect(slide2, x, top_y, box_w, box_h,
             fill_color=LIGHT_GRAY if i % 2 == 0 else WHITE,
             line_color=ORANGE, line_width=Pt(2))

    # 번호 원
    circle = slide2.shapes.add_shape(9, x + Inches(0.8), top_y + Inches(0.2),
                                      Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = num; run.font.size = Pt(16); run.font.bold = True
    run.font.color.rgb = WHITE

    # 제목
    add_text(slide2, title, x + Inches(0.1), top_y + Inches(0.85), box_w - Inches(0.2), Inches(0.5),
             font_size=Pt(14), bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # 설명
    add_text(slide2, desc, x + Inches(0.1), top_y + Inches(1.5), box_w - Inches(0.2), Inches(1.2),
             font_size=Pt(12), color=GRAY, align=PP_ALIGN.CENTER)

    # 화살표 (마지막 제외)
    if i < 4:
        arr_x = x + box_w + Inches(0.02)
        add_text(slide2, "▶", arr_x, top_y + box_h/2 - Pt(15),
                 gap, Pt(30), font_size=Pt(16), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# 주석
add_text(slide2, "※ 결제 수단: 신용카드 / 체크카드 / 카카오페이 / 네이버페이",
         Inches(0.5), Inches(5.3), W - Inches(1), Pt(25),
         font_size=Pt(11), color=GRAY, italic=True)


# ════════════════════════════════════════════════════════════
# 슬라이드 3: STEP 1 — 서비스 접속
# ════════════════════════════════════════════════════════════
slide3 = blank_slide(prs)
add_rect(slide3, 0, 0, W, Inches(1.1), fill_color=ORANGE)
add_text(slide3, "STEP 1  |  서비스 접속", Inches(0.5), Inches(0.2), W, Inches(0.7),
         font_size=Pt(22), bold=True, color=WHITE)

# 좌측: 설명
add_text(slide3, "① 브라우저 또는 앱에서 접속", Inches(0.5), Inches(1.4), Inches(5.5), Pt(28),
         font_size=Pt(15), bold=True, color=DARK)
add_text(slide3, "https://sourcing-kit.kr", Inches(0.5), Inches(1.9), Inches(5.5), Pt(25),
         font_size=Pt(13), color=TOSS_BLUE, italic=True)

add_text(slide3, "소싱킷은 무역 소싱 관리 SaaS 서비스로,\n"
                 "웹 브라우저 및 PWA(앱 설치) 형태로\n"
                 "서비스됩니다.\n\n"
                 "• URL 직접 접속 가능\n"
                 "• Google / 카카오 소셜 로그인\n"
                 "• 모바일/PC 반응형 지원",
         Inches(0.5), Inches(2.5), Inches(5.5), Inches(3),
         font_size=Pt(13), color=GRAY)

# 우측: 화면 시뮬레이션
phone_l = Inches(7.2); phone_t = Inches(1.2); phone_w = Inches(5.0); phone_h = Inches(5.8)
add_rect(slide3, phone_l, phone_t, phone_w, phone_h,
         fill_color=WHITE, line_color=ORANGE, line_width=Pt(2))
add_rect(slide3, phone_l, phone_t, phone_w, Inches(1.0), fill_color=ORANGE)
add_text(slide3, "소싱킷 | 무역 소싱 관리", phone_l + Inches(0.2), phone_t + Inches(0.2),
         phone_w - Inches(0.4), Pt(30), font_size=Pt(13), bold=True, color=WHITE)
add_text(slide3, "🔍  소싱킷", phone_l + Inches(0.2), phone_t + Inches(1.2),
         phone_w - Inches(0.4), Pt(35), font_size=Pt(18), bold=True, color=DARK)
add_text(slide3, "이우 무역상을 위한 소싱 원가계산 앱",
         phone_l + Inches(0.2), phone_t + Inches(1.8), phone_w - Inches(0.4), Pt(25),
         font_size=Pt(11), color=GRAY)
add_rect(slide3, phone_l + Inches(0.5), phone_t + Inches(2.4),
         phone_w - Inches(1), Inches(0.55), fill_color=ORANGE)
add_text(slide3, "시작하기 →", phone_l + Inches(0.5), phone_t + Inches(2.4),
         phone_w - Inches(1), Inches(0.55), font_size=Pt(14), bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 슬라이드 4: STEP 2 — 플랜 선택 (가격 페이지)
# ════════════════════════════════════════════════════════════
slide4 = blank_slide(prs)
add_rect(slide4, 0, 0, W, Inches(1.1), fill_color=ORANGE)
add_text(slide4, "STEP 2  |  가격 페이지 — 플랜 선택", Inches(0.5), Inches(0.2), W, Inches(0.7),
         font_size=Pt(22), bold=True, color=WHITE)

# 가격 카드 3개
cards = [
    ("무료", "₩0", "무료 플랜", ["상품 10개", "공급업체 5곳", "AI 분석 3회/일"], GRAY, False),
    ("맛보기", "₩9,900", "30일 이용권", ["상품 100개", "공급업체 30곳", "AI 분석 20회/일", "엑셀/PDF 내보내기"], RGBColor(0xD9, 0x77, 0x06), False),
    ("Pro", "₩7,900", "/월 자동결제", ["상품 무제한", "공급업체 무제한", "AI 분석 무제한", "엑셀/PDF 내보내기", "팀 공유"], ORANGE, True),
]

card_w = Inches(3.8)
card_h = Inches(5.5)
start_x2 = Inches(0.5)
for i, (name, price, sub, feats, color, highlight) in enumerate(cards):
    x = start_x2 + (card_w + Inches(0.35)) * i
    line_color = ORANGE if highlight else RGBColor(0xD1, 0xD5, 0xDB)
    add_rect(slide4, x, Inches(1.3), card_w, card_h,
             fill_color=WHITE, line_color=line_color, line_width=Pt(3 if highlight else 1.5))
    add_text(slide4, name, x + Inches(0.2), Inches(1.5), card_w - Inches(0.4), Pt(30),
             font_size=Pt(16), bold=True, color=color)
    add_text(slide4, price, x + Inches(0.2), Inches(2.0), card_w - Inches(0.4), Pt(35),
             font_size=Pt(22), bold=True, color=DARK)
    add_text(slide4, sub, x + Inches(0.2), Inches(2.6), card_w - Inches(0.4), Pt(22),
             font_size=Pt(11), color=GRAY)
    for j, feat in enumerate(feats):
        add_text(slide4, f"✓  {feat}", x + Inches(0.2), Inches(3.1) + Pt(22)*j,
                 card_w - Inches(0.4), Pt(22), font_size=Pt(11), color=DARK)
    btn_color = ORANGE if highlight else RGBColor(0xD9, 0x77, 0x06) if name == "맛보기" else LIGHT_GRAY
    btn_text_color = WHITE if highlight or name == "맛보기" else GRAY
    add_rect(slide4, x + Inches(0.2), Inches(6.3), card_w - Inches(0.4), Inches(0.35),
             fill_color=btn_color)
    label = "무료 시작" if name == "무료" else ("30일 맛보기 시작" if name == "맛보기" else "⚡ Pro 구독하기")
    add_text(slide4, label, x + Inches(0.2), Inches(6.3), card_w - Inches(0.4), Inches(0.35),
             font_size=Pt(11), bold=True, color=btn_text_color, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 슬라이드 5: STEP 3 — 결제하기 버튼 클릭
# ════════════════════════════════════════════════════════════
slide5 = blank_slide(prs)
add_rect(slide5, 0, 0, W, Inches(1.1), fill_color=ORANGE)
add_text(slide5, "STEP 3  |  결제하기 버튼 클릭", Inches(0.5), Inches(0.2), W, Inches(0.7),
         font_size=Pt(22), bold=True, color=WHITE)

# 좌측 설명
add_text(slide5, "맛보기 (₩9,900 · 단건결제)", Inches(0.5), Inches(1.5), Inches(5.5), Pt(28),
         font_size=Pt(16), bold=True, color=RGBColor(0xD9, 0x77, 0x06))
add_text(slide5, "• 1회 결제, 자동갱신 없음\n• 30일 이용 후 자동 종료\n• 결제 즉시 이용 가능",
         Inches(0.5), Inches(2.1), Inches(5.5), Inches(1.5), font_size=Pt(13), color=GRAY)

add_text(slide5, "Pro 구독 (₩7,900/월 또는 ₩70,800/년)", Inches(0.5), Inches(3.8), Inches(5.5), Pt(28),
         font_size=Pt(16), bold=True, color=ORANGE)
add_text(slide5, "• 카드 빌링키 발급 후 자동결제\n• 매월 자동 청구 (언제든 취소 가능)\n• 연결제 시 25% 할인",
         Inches(0.5), Inches(4.4), Inches(5.5), Inches(1.5), font_size=Pt(13), color=GRAY)

# 우측 버튼 예시
add_rect(slide5, Inches(7), Inches(1.5), Inches(5.5), Inches(1.5),
         fill_color=RGBColor(0xFE, 0xF3, 0xC7), line_color=RGBColor(0xD9, 0x77, 0x06), line_width=Pt(2))
add_text(slide5, "30일 맛보기 시작 ₩9,900",
         Inches(7.1), Inches(1.65), Inches(5.3), Inches(0.8),
         font_size=Pt(17), bold=True, color=RGBColor(0xD9, 0x77, 0x06), align=PP_ALIGN.CENTER)
add_text(slide5, "↑  이 버튼 클릭 시 토스페이먼츠 결제창 오픈",
         Inches(7), Inches(3.1), Inches(5.5), Pt(25),
         font_size=Pt(11), color=GRAY, italic=True, align=PP_ALIGN.CENTER)

add_rect(slide5, Inches(7), Inches(4.0), Inches(5.5), Inches(1.5),
         fill_color=RGBColor(0xFF, 0xF3, 0xEA), line_color=ORANGE, line_width=Pt(2))
add_text(slide5, "⚡ Pro 구독하기 ₩7,900/월",
         Inches(7.1), Inches(4.15), Inches(5.3), Inches(0.8),
         font_size=Pt(17), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(slide5, "↑  이 버튼 클릭 시 빌링키 발급 결제창 오픈",
         Inches(7), Inches(5.6), Inches(5.5), Pt(25),
         font_size=Pt(11), color=GRAY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 슬라이드 6: STEP 4 — 토스페이먼츠 결제창
# ════════════════════════════════════════════════════════════
slide6 = blank_slide(prs)
add_rect(slide6, 0, 0, W, Inches(1.1), fill_color=ORANGE)
add_text(slide6, "STEP 4  |  토스페이먼츠 결제창", Inches(0.5), Inches(0.2), W, Inches(0.7),
         font_size=Pt(22), bold=True, color=WHITE)

# 좌측 설명
add_text(slide6, "결제 방법", Inches(0.5), Inches(1.4), Inches(5), Pt(28),
         font_size=Pt(15), bold=True, color=DARK)
methods = ["💳  신용카드 / 체크카드", "🟡  카카오페이", "🟢  네이버페이"]
for i, m in enumerate(methods):
    add_text(slide6, m, Inches(0.8), Inches(2.0) + Pt(28)*i, Inches(5), Pt(28),
             font_size=Pt(13), color=DARK)

add_text(slide6, "빌링키 결제 흐름 (Pro 구독)", Inches(0.5), Inches(3.5), Inches(5), Pt(28),
         font_size=Pt(14), bold=True, color=DARK)
add_text(slide6, "카드 정보 입력 → 빌링키 발급\n→ 최초 승인 → 매월 자동청구",
         Inches(0.5), Inches(4.1), Inches(5), Inches(1.0), font_size=Pt(12), color=GRAY)

add_text(slide6, "단건 결제 흐름 (맛보기)", Inches(0.5), Inches(5.3), Inches(5), Pt(28),
         font_size=Pt(14), bold=True, color=DARK)
add_text(slide6, "결제 수단 선택 → 인증 → 즉시 승인",
         Inches(0.5), Inches(5.9), Inches(5), Inches(0.6), font_size=Pt(12), color=GRAY)

# 우측: 결제창 시뮬레이션
pw_l = Inches(6.5); pw_t = Inches(1.3); pw_w = Inches(6.3); pw_h = Inches(5.8)
add_rect(slide6, pw_l, pw_t, pw_w, pw_h, fill_color=WHITE,
         line_color=RGBColor(0xE5, 0xE7, 0xEB), line_width=Pt(1.5))
# Toss 헤더
add_rect(slide6, pw_l, pw_t, pw_w, Inches(0.7), fill_color=TOSS_BLUE)
add_text(slide6, "토스페이먼츠", pw_l + Inches(0.2), pw_t + Inches(0.1),
         pw_w - Inches(0.4), Pt(30), font_size=Pt(16), bold=True, color=WHITE)
# 주문 정보
add_text(slide6, "주문 정보", pw_l + Inches(0.3), pw_t + Inches(0.9),
         pw_w - Inches(0.5), Pt(22), font_size=Pt(11), bold=True, color=DARK)
add_text(slide6, "소싱킷 Pro 구독  |  ₩7,900",
         pw_l + Inches(0.3), pw_t + Inches(1.25), pw_w - Inches(0.5), Pt(22),
         font_size=Pt(12), color=DARK)
# 결제 수단 탭
for j, (method, color) in enumerate([("카드", TOSS_BLUE), ("카카오페이", RGBColor(0xFE, 0xE5, 0x00)), ("네이버페이", RGBColor(0x03, 0xC7, 0x5A))]):
    mx = pw_l + Inches(0.3) + Inches(1.9) * j
    add_rect(slide6, mx, pw_t + Inches(1.8), Inches(1.8), Inches(0.45),
             fill_color=color, line_color=None)
    tc = WHITE if method != "카카오페이" else DARK
    add_text(slide6, method, mx, pw_t + Inches(1.8), Inches(1.8), Inches(0.45),
             font_size=Pt(12), bold=True, color=tc, align=PP_ALIGN.CENTER)
# 카드 번호 필드
add_rect(slide6, pw_l + Inches(0.3), pw_t + Inches(2.55), pw_w - Inches(0.5), Inches(0.5),
         fill_color=LIGHT_GRAY, line_color=RGBColor(0xD1, 0xD5, 0xDB), line_width=Pt(1))
add_text(slide6, "카드 번호  0000 - 0000 - 0000 - 0000",
         pw_l + Inches(0.5), pw_t + Inches(2.6), pw_w - Inches(0.8), Pt(28),
         font_size=Pt(11), color=GRAY)
# 유효기간 / CVC
add_rect(slide6, pw_l + Inches(0.3), pw_t + Inches(3.2), Inches(2.5), Inches(0.5),
         fill_color=LIGHT_GRAY, line_color=RGBColor(0xD1, 0xD5, 0xDB), line_width=Pt(1))
add_text(slide6, "유효기간  MM / YY",
         pw_l + Inches(0.5), pw_t + Inches(3.25), Inches(2.3), Pt(28), font_size=Pt(11), color=GRAY)
add_rect(slide6, pw_l + Inches(3.0), pw_t + Inches(3.2), Inches(2.8), Inches(0.5),
         fill_color=LIGHT_GRAY, line_color=RGBColor(0xD1, 0xD5, 0xDB), line_width=Pt(1))
add_text(slide6, "CVC  000",
         pw_l + Inches(3.2), pw_t + Inches(3.25), Inches(2.6), Pt(28), font_size=Pt(11), color=GRAY)
# 결제 버튼
add_rect(slide6, pw_l + Inches(0.3), pw_t + Inches(4.6), pw_w - Inches(0.5), Inches(0.6),
         fill_color=TOSS_BLUE)
add_text(slide6, "₩7,900  결제하기",
         pw_l + Inches(0.3), pw_t + Inches(4.65), pw_w - Inches(0.5), Inches(0.55),
         font_size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 슬라이드 7: STEP 5 — 결제 완료 & 서비스 이용
# ════════════════════════════════════════════════════════════
slide7 = blank_slide(prs)
add_rect(slide7, 0, 0, W, Inches(1.1), fill_color=ORANGE)
add_text(slide7, "STEP 5  |  결제 완료 및 서비스 이용", Inches(0.5), Inches(0.2), W, Inches(0.7),
         font_size=Pt(22), bold=True, color=WHITE)

# 완료 배지
add_rect(slide7, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.2),
         fill_color=RGBColor(0xD1, 0xFA, 0xE5), line_color=GREEN, line_width=Pt(2))
add_text(slide7, "✅  결제 승인 완료  —  서비스가 즉시 활성화됩니다",
         Inches(0.7), Inches(1.45), Inches(12), Inches(0.9),
         font_size=Pt(18), bold=True, color=GREEN)

# 결제 후 화면 2개
# 왼쪽: 결제 완료 화면
add_rect(slide7, Inches(0.5), Inches(2.7), Inches(5.8), Inches(4.0),
         fill_color=WHITE, line_color=RGBColor(0xD1, 0xD5, 0xDB), line_width=Pt(1.5))
add_rect(slide7, Inches(0.5), Inches(2.7), Inches(5.8), Inches(0.6), fill_color=GREEN)
add_text(slide7, "결제 완료", Inches(0.7), Inches(2.75), Inches(5.4), Pt(30),
         font_size=Pt(14), bold=True, color=WHITE)
add_text(slide7, "🎉", Inches(2.5), Inches(3.4), Inches(1.8), Inches(0.8), font_size=Pt(36), align=PP_ALIGN.CENTER)
add_text(slide7, "소싱킷 Pro 플랜\n활성화 완료!",
         Inches(0.7), Inches(4.3), Inches(5.4), Inches(0.9),
         font_size=Pt(16), bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(slide7, "결제일 2025.01.01  |  다음 결제일 2025.02.01",
         Inches(0.7), Inches(5.3), Inches(5.4), Pt(22),
         font_size=Pt(10), color=GRAY, align=PP_ALIGN.CENTER)
add_rect(slide7, Inches(1.0), Inches(5.8), Inches(4.8), Inches(0.5), fill_color=ORANGE)
add_text(slide7, "서비스 시작하기 →",
         Inches(1.0), Inches(5.82), Inches(4.8), Pt(30),
         font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 오른쪽: 구독 관리 화면
add_rect(slide7, Inches(7.0), Inches(2.7), Inches(5.8), Inches(4.0),
         fill_color=WHITE, line_color=RGBColor(0xD1, 0xD5, 0xDB), line_width=Pt(1.5))
add_rect(slide7, Inches(7.0), Inches(2.7), Inches(5.8), Inches(0.6), fill_color=ORANGE)
add_text(slide7, "구독 관리", Inches(7.2), Inches(2.75), Inches(5.4), Pt(30),
         font_size=Pt(14), bold=True, color=WHITE)
rows = [("현재 플랜", "Pro"), ("상태", "이용 중"), ("결제 금액", "₩7,900/월"), ("다음 결제일", "2025.02.01")]
for k, (label, val) in enumerate(rows):
    y = Inches(3.5) + Inches(0.5) * k
    add_text(slide7, label, Inches(7.2), y, Inches(2.5), Pt(22), font_size=Pt(11), color=GRAY)
    add_text(slide7, val, Inches(10.0), y, Inches(2.5), Pt(22),
             font_size=Pt(12), bold=True, color=DARK, align=PP_ALIGN.RIGHT)
add_text(slide7, "* 구독 관리 화면에서 언제든지 취소 가능",
         Inches(7.2), Inches(5.8), Inches(5.4), Pt(22), font_size=Pt(10), color=GRAY, italic=True)


# ════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════
out_path = r"C:\Users\USER\Documents\Claude\cord\trade-sourcing-app\scripts\소싱킷_결제경로.pptx"
prs.save(out_path)
print(f"PPT 저장 완료: {out_path}")
print(f"슬라이드 수: {len(prs.slides)}장")
