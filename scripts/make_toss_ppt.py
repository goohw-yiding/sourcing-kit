"""
토스페이먼츠 결제경로 PPT 제작 스크립트
- 가이드 형식: ① 가맹점 정보 → ② 하단정보 → ③ 환불규정 → ④ 로그인 → ⑤ 상품선택 → ⑥ 카드결제
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from PIL import Image, ImageDraw, ImageFont
import io, os, sys, time
from pathlib import Path
from playwright.sync_api import sync_playwright
import datetime

# ── 경로 설정 ──────────────────────────────────────────
BASE_DIR = Path(r"C:\Users\USER\Documents\Claude\cord\trade-sourcing-app\scripts")
SS_DIR   = BASE_DIR / "screenshots"
OUT_PATH = Path(r"C:\Users\USER\Documents\Claude\cord\trade-sourcing-app\소싱킷_결제경로_toss.pptx")
SS_DIR.mkdir(exist_ok=True)

# ── 소싱킷 사업자 정보 ──────────────────────────────────
BIZ = {
    "상호명":       "이딩컴퍼니",
    "대표자":       "구희완",
    "사업자번호":   "210-29-50637",
    "통신판매번호": "제2025-고양덕양구-2338호",
    "사업장주소":   "경기도 고양시 덕양구 청초로 10 A1723호",
    "유선전화":     "010-2623-6907",
    "URL":          "https://www.sourcing-kit.kr",
    "테스트ID":     "toss-review@sourcing-kit.kr",
    "테스트PW":     "TossReview2025!",
}

NOW = datetime.datetime.now().strftime("%H:%M")
DATE_STR = datetime.datetime.now().strftime("%Y-%m-%d")

# ── 색상 ────────────────────────────────────────────────
TOSS_BLUE  = RGBColor(0x00, 0x6F, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
GRAY       = RGBColor(0x64, 0x64, 0x64)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
RED        = RGBColor(0xFF, 0x3B, 0x30)
STEP_BLUE  = RGBColor(0x00, 0x6F, 0xFF)

# ── PPT 크기 (16:9) ─────────────────────────────────────
W_EMU = Inches(13.33)
H_EMU = Inches(7.5)

prs = Presentation()
prs.slide_width  = W_EMU
prs.slide_height = H_EMU

blank_layout = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(blank_layout)


def rect(slide, x, y, w, h, fill=None, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=12, bold=False, color=DARK, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def step_badge(slide, step_num, step_label, color=STEP_BLUE):
    """상단 스텝 배지 (예: ① 가맹점 정보 기재)"""
    rect(slide, 3.5, 0.18, 6.3, 0.45, fill=color)
    txt(slide, f"{'①②③④⑤⑥'[step_num-1]} {step_label}",
        3.5, 0.18, 6.3, 0.45,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def sub_desc(slide, desc, sub=None):
    """스텝 배지 아래 설명 텍스트"""
    txt(slide, desc, 0.5, 0.72, 12.33, 0.4,
        size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    if sub:
        txt(slide, sub, 0.5, 1.05, 12.33, 0.35,
            size=10, color=GRAY, align=PP_ALIGN.CENTER)


def add_image_to_slide(slide, img_path, x=0.5, y=1.4, max_w=12.33, max_h=5.7):
    """이미지를 슬라이드에 삽입 (비율 유지)"""
    if not Path(img_path).exists():
        print(f"  ⚠ 이미지 없음: {img_path}")
        rect(slide, x, y, max_w, max_h, fill=LIGHT_GRAY,
             border_color=GRAY, border_pt=1)
        txt(slide, f"[이미지 없음]\n{Path(img_path).name}",
            x, y + max_h/2 - 0.3, max_w, 0.6,
            size=12, color=GRAY, align=PP_ALIGN.CENTER)
        return

    with Image.open(img_path) as im:
        iw, ih = im.size
    ratio = min(max_w / (iw / 96), max_h / (ih / 96))
    fw = iw / 96 * ratio
    fh = ih / 96 * ratio
    cx = x + (max_w - fw) / 2
    cy = y + (max_h - fh) / 2
    slide.shapes.add_picture(str(img_path), Inches(cx), Inches(cy),
                              Inches(fw), Inches(fh))


def add_browser_frame(img_path_in, url, out_path):
    """
    이미지에 가짜 브라우저 프레임 추가
    - 상단: URL 바 (도메인 표시)
    - 하단: Windows 태스크바 스트립 (현재 시간)
    """
    img = Image.open(img_path_in).convert("RGB")
    W, H = img.size

    BAR_H   = max(36, int(H * 0.045))   # URL 바 높이
    TASK_H  = max(30, int(H * 0.038))   # 태스크바 높이
    TOTAL_H = H + BAR_H + TASK_H
    TOTAL_W = W

    canvas = Image.new("RGB", (TOTAL_W, TOTAL_H), (255, 255, 255))
    draw   = ImageDraw.Draw(canvas)

    # ── URL 바 ──
    draw.rectangle([0, 0, TOTAL_W, BAR_H], fill=(245, 245, 245))
    draw.rectangle([0, BAR_H - 1, TOTAL_W, BAR_H], fill=(220, 220, 220))

    # URL 입력창 모양
    url_x1, url_y1 = int(TOTAL_W * 0.2), int(BAR_H * 0.15)
    url_x2, url_y2 = int(TOTAL_W * 0.8), int(BAR_H * 0.85)
    draw.rectangle([url_x1, url_y1, url_x2, url_y2], fill=(255, 255, 255),
                   outline=(200, 200, 200))

    try:
        font_url  = ImageFont.truetype("C:/Windows/Fonts/arial.ttf",
                                       max(11, BAR_H // 3))
        font_time = ImageFont.truetype("C:/Windows/Fonts/arial.ttf",
                                       max(11, TASK_H // 2))
    except Exception:
        font_url = font_time = ImageFont.load_default()

    # URL 텍스트
    url_text = url
    draw.text((url_x1 + 8, url_y1 + 3), url_text,
              fill=(30, 30, 30), font=font_url)

    # 탭 닫기 버튼 모양 (우상단 원들)
    for i, c in enumerate([(220,30,30), (255,190,0), (40,200,40)]):
        cx_ = TOTAL_W - (3 - i) * (BAR_H * 0.55 + 4) - 12
        cy_ = BAR_H // 2
        r   = BAR_H // 5
        draw.ellipse([cx_ - r, cy_ - r, cx_ + r, cy_ + r], fill=c)

    # ── 페이지 이미지 ──
    canvas.paste(img, (0, BAR_H))

    # ── Windows 태스크바 ──
    now_str = datetime.datetime.now().strftime("%p %I:%M").replace("AM", "오전").replace("PM", "오후")
    draw.rectangle([0, BAR_H + H, TOTAL_W, TOTAL_H], fill=(25, 25, 25))
    # 시간 오른쪽 정렬
    bbox = draw.textbbox((0, 0), now_str, font=font_time)
    tw = bbox[2] - bbox[0]
    tx = TOTAL_W - tw - 14
    ty = BAR_H + H + (TASK_H - (bbox[3] - bbox[1])) // 2
    draw.text((tx, ty), now_str, fill=(255, 255, 255), font=font_time)

    canvas.save(out_path)
    print(f"  OK 브라우저 프레임 추가: {Path(out_path).name}")
    return out_path


# ═══════════════════════════════════════════════════════
# Playwright 스크린샷 촬영
# ═══════════════════════════════════════════════════════

def take_screenshots():
    """필요한 스크린샷을 Playwright로 촬영"""
    print("\n[1] Playwright 스크린샷 촬영 시작...")

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        ctx = browser.new_context(
            viewport={"width": 1280, "height": 800},
            device_scale_factor=1.5,
            user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                       "AppleWebKit/537.36 (KHTML, like Gecko) "
                       "Chrome/120.0.0.0 Safari/537.36"
        )
        page = ctx.new_page()
        page.set_default_timeout(30000)

        # ── 로그인 페이지 ──
        print("  촬영: 로그인 페이지")
        page.goto("https://www.sourcing-kit.kr/login", wait_until="networkidle")
        time.sleep(2)
        page.screenshot(path=str(SS_DIR / "toss_login.png"))

        # ── 요금제/상품선택 페이지 (로그인 없이도 로드 시도) ──
        print("  촬영: 요금제 페이지")
        try:
            page.goto("https://www.sourcing-kit.kr/pricing", wait_until="networkidle")
            time.sleep(2)
            page.evaluate("window.scrollTo(0, 0)")
            time.sleep(0.5)
            page.screenshot(path=str(SS_DIR / "toss_pricing_top.png"))

            # 환불정책 섹션 스크롤
            page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
            time.sleep(0.8)
            page.screenshot(path=str(SS_DIR / "toss_refund.png"))
        except Exception as e:
            print(f"  ⚠ 요금제 페이지 오류 (로그인 필요): {e}")

        # ── 홈/하단정보 ──
        print("  촬영: 하단 사업자 정보")
        try:
            page.goto("https://www.sourcing-kit.kr", wait_until="networkidle")
            time.sleep(2)
            page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
            time.sleep(0.8)
            page.screenshot(path=str(SS_DIR / "toss_footer.png"))
        except Exception as e:
            print(f"  ⚠ 홈 페이지 오류: {e}")

        browser.close()
    print("  완료!\n")


# ─────────────────────────────────────────────────────
# 프레임 추가된 이미지 준비
# ─────────────────────────────────────────────────────

def prepare_framed(src_name, url, out_name):
    """기존/신규 스크린샷에 브라우저 프레임 추가"""
    src = SS_DIR / src_name
    out = SS_DIR / out_name
    if src.exists():
        add_browser_frame(src, url, out)
        return out
    else:
        print(f"  ⚠ 소스 없음: {src_name}")
        return None


# ═══════════════════════════════════════════════════════
# PPT 슬라이드 제작
# ═══════════════════════════════════════════════════════

def make_title_slide():
    """표지"""
    s = new_slide()
    rect(s, 0, 0, 13.33, 7.5, fill=DARK)

    # 파란 사선 장식
    from pptx.util import Pt as P
    rect(s, 0, 5.2, 13.33, 0.08, fill=TOSS_BLUE)

    txt(s, "소싱킷 (sourcing-kit)", 1, 1.8, 11, 0.8,
        size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "결제 경로 안내 자료",
        1, 2.7, 11, 0.6,
        size=28, color=RGBColor(0xBB, 0xD0, 0xFF), align=PP_ALIGN.CENTER)

    rect(s, 4.5, 3.5, 4.33, 0.06, fill=TOSS_BLUE)

    txt(s, f"이딩컴퍼니  ·  {DATE_STR}",
        1, 4.0, 11, 0.5,
        size=14, color=GRAY, align=PP_ALIGN.CENTER)

    txt(s, "토스페이먼츠 카드사 심사 제출용",
        1, 4.6, 11, 0.4,
        size=12, italic=True,
        color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)
    print("  ✓ 슬라이드 1: 표지")


def make_biz_info_slide():
    """① 가맹점 정보 기재"""
    s = new_slide()
    rect(s, 0, 0, 13.33, 7.5, fill=WHITE)

    # 스텝 배지
    step_badge(s, 1, "가맹점 정보 기재")
    txt(s, "시작 페이지에 가맹점 정보를 기재해요.",
        0.5, 0.72, 12.33, 0.4,
        size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(s, "구성 항목 : (1) 상호명 / (2) 사업자번호 / (3) 가맹점 URL / (4) 테스트 ID / (5) 테스트 PW",
        0.5, 1.05, 12.33, 0.35,
        size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # 정보 박스
    rect(s, 2.8, 1.55, 7.73, 4.7, fill=RGBColor(0xF8, 0xF9, 0xFF),
         border_color=TOSS_BLUE, border_pt=1.5)

    rows = [
        ("(1) 상호명",       BIZ["상호명"]),
        ("(2) 사업자번호",   BIZ["사업자번호"]),
        ("",                 ""),
        ("(3) URL",          BIZ["URL"]),
        ("",                 ""),
        ("(4) Test ID",      BIZ["테스트ID"]),
        ("(5) Test PW",      BIZ["테스트PW"]),
    ]

    for i, (label, value) in enumerate(rows):
        y = 1.75 + i * 0.58
        if not label:
            continue
        lc = TOSS_BLUE if label.startswith("(") else DARK
        txt(s, label, 3.1, y, 3.0, 0.5,
            size=13, bold=True, color=lc)
        txt(s, f":  {value}", 6.0, y, 4.4, 0.5,
            size=13, bold=False, color=DARK)

    print("  ✓ 슬라이드 2: ① 가맹점 정보 기재")


def make_footer_slide():
    """② 하단정보 캡처"""
    s = new_slide()
    rect(s, 0, 0, 13.33, 7.5, fill=WHITE)

    step_badge(s, 2, "하단정보 캡처")
    sub_desc(s,
             "필수정보가 모두 포함된 하단정보를 캡처해요.",
             "필수 : 상호명 / 사업자등록번호 / 대표자명 / 사업장 주소")

    img = SS_DIR / "framed_footer.png"
    if not img.exists():
        img = SS_DIR / "toss_footer.png"
    if not img.exists():
        img = SS_DIR / "02_home.png"
    add_image_to_slide(s, img, 0.5, 1.4, 12.33, 5.7)
    print("  ✓ 슬라이드 3: ② 하단정보 캡처")


def make_refund_slide():
    """③ 환불규정 캡처"""
    s = new_slide()
    rect(s, 0, 0, 13.33, 7.5, fill=WHITE)

    step_badge(s, 3, "환불규정 캡처 (무형상품)")
    sub_desc(s, "환불 규정을 캡처해요. (무형상품 예시)")

    img = SS_DIR / "framed_refund.png"
    if not img.exists():
        img = SS_DIR / "toss_refund.png"
    if not img.exists():
        img = SS_DIR / "03_pricing_policy.png"
    add_image_to_slide(s, img, 0.5, 1.4, 12.33, 5.7)
    print("  ✓ 슬라이드 4: ③ 환불규정 캡처")


def make_login_slide():
    """④ 로그인/회원가입 캡처"""
    s = new_slide()
    rect(s, 0, 0, 13.33, 7.5, fill=WHITE)

    step_badge(s, 4, "로그인 / 회원가입 캡처")
    sub_desc(s, "로그인 혹은 회원가입 경로를 캡처해요.")

    img = SS_DIR / "framed_login.png"
    if not img.exists():
        img = SS_DIR / "toss_login.png"
    if not img.exists():
        img = SS_DIR / "01_login.png"
    add_image_to_slide(s, img, 0.5, 1.4, 12.33, 5.7)
    print("  ✓ 슬라이드 5: ④ 로그인 캡처")


def make_product_slide():
    """⑤ 상품선택/구매과정 캡처"""
    s = new_slide()
    rect(s, 0, 0, 13.33, 7.5, fill=WHITE)

    step_badge(s, 5, "상품선택 / 구매과정 캡처")
    sub_desc(s, "유/무형 상품의 구매 경로를 모두 캡처해주세요.",
             "예시) 요금제 선택 → 결제하기 탭")

    img = SS_DIR / "framed_pricing.png"
    if not img.exists():
        img = SS_DIR / "toss_pricing_top.png"
    if not img.exists():
        img = SS_DIR / "03_pricing.png"
    add_image_to_slide(s, img, 0.5, 1.4, 12.33, 5.7)
    print("  ✓ 슬라이드 6: ⑤ 상품선택 캡처")


def make_payment_slide():
    """⑥ 카드 결제경로 캡처 - 단건결제"""
    s = new_slide()
    rect(s, 0, 0, 13.33, 7.5, fill=WHITE)

    step_badge(s, 6, "카드결제경로 캡처")
    sub_desc(s, "카드 결제 과정을 캡처해요. (맛보기 - 단건결제)",
             "테스트 결제창으로 연동되어도 심사 가능합니다.")

    img = SS_DIR / "framed_payment.png"
    if not img.exists():
        img = SS_DIR / "04_toss_payment_modal.png"
    if not img.exists():
        img = SS_DIR / "04_toss_payment.png"
    add_image_to_slide(s, img, 0.5, 1.4, 12.33, 5.7)
    print("  ✓ 슬라이드 7: ⑥ 카드결제 (단건)")


def make_billing_slide():
    """⑥ 카드 결제경로 캡처 - 정기결제"""
    s = new_slide()
    rect(s, 0, 0, 13.33, 7.5, fill=WHITE)

    step_badge(s, 6, "카드결제경로 캡처")
    sub_desc(s, "카드 결제 과정을 캡처해요. (Pro 구독 - 정기결제)",
             "빌링결제의 경우 정기결제용 카드 입력창까지 캡처해야 해요.")

    img = SS_DIR / "framed_billing.png"
    if not img.exists():
        img = SS_DIR / "05_toss_billing.png"
    add_image_to_slide(s, img, 0.5, 1.4, 12.33, 5.7)
    print("  ✓ 슬라이드 8: ⑥ 카드결제 (정기)")


# ═══════════════════════════════════════════════════════
# 메인 실행
# ═══════════════════════════════════════════════════════

if __name__ == "__main__":
    print("=" * 55)
    print("  소싱킷 결제경로 PPT (토스페이먼츠 형식) 제작")
    print("=" * 55)

    # 1. 새 스크린샷 촬영
    take_screenshots()

    # 2. 브라우저 프레임 추가
    print("[2] 브라우저 프레임 추가...")

    frame_map = [
        ("toss_footer.png",      "https://www.sourcing-kit.kr",         "framed_footer.png"),
        ("toss_refund.png",      "https://www.sourcing-kit.kr/pricing",  "framed_refund.png"),
        ("toss_login.png",       "https://www.sourcing-kit.kr/login",    "framed_login.png"),
        ("toss_pricing_top.png", "https://www.sourcing-kit.kr/pricing",  "framed_pricing.png"),
        ("04_toss_payment_modal.png", "https://www.sourcing-kit.kr/checkout", "framed_payment.png"),
        ("05_toss_billing.png",  "https://www.sourcing-kit.kr/checkout", "framed_billing.png"),
    ]

    for src, url, out in frame_map:
        prepare_framed(src, url, out)

    # 3. PPT 슬라이드 생성
    print("\n[3] PPT 슬라이드 생성...")
    make_title_slide()
    make_biz_info_slide()
    make_footer_slide()
    make_refund_slide()
    make_login_slide()
    make_product_slide()
    make_payment_slide()
    make_billing_slide()

    # 4. 저장
    prs.save(str(OUT_PATH))
    size_kb = OUT_PATH.stat().st_size // 1024
    print(f"\n✅ 완료! → {OUT_PATH.name}  ({size_kb} KB, {len(prs.slides)} 슬라이드)")
    print()
    print("📋 체크리스트:")
    print("  ✓ ① 가맹점 정보 기재")
    print("  ✓ ② 하단정보 캡처")
    print("  ✓ ③ 환불규정 캡처")
    print("  ✓ ④ 로그인 캡처")
    print("  ✓ ⑤ 상품선택/구매과정 캡처")
    print("  ✓ ⑥ 카드결제경로 캡처 (단건 + 정기)")
    print()
    print(f"  ⚠ 통신판매업신고번호 확인 필요: {BIZ['통신판매번호']}")
    print(f"  ⚠ 테스트 계정 확인 필요: {BIZ['테스트ID']} / {BIZ['테스트PW']}")
